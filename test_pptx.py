import collections
import collections.abc
from pptx import Presentation

def extract_template(input_pptx, output_pptx):
    prs = Presentation(input_pptx)
    # Remove all slides to keep only the master and layouts
    # We have to delete them from XML since python-pptx doesn't have a direct prs.slides.remove() or similar
    
    # Or in python-pptx, deleting slides from prs:
    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    for slide in slides:
        xml_slides.remove(slide)
        
    prs.save(output_pptx)
    print("Saved template with 0 slides.")

