from pptx import Presentation
import test_pptx

# Create a mock presentation with 1 slide to test
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
prs.save("mock.pptx")

test_pptx.extract_template("mock.pptx", "mock_template.pptx")

out_prs = Presentation("mock_template.pptx")
print("Number of slides after extraction:", len(out_prs.slides))
