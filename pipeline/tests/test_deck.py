"""Tests for deck builders: one_pager and implementation_plan.

Both builders need only a Mission, a MissionResult, a viz dict, an
EntitySchema, and a writable output path.  No DB required.

Covers:
  - File is created on disk
  - File is a valid PPTX (operable with python-pptx)
  - Missing kaleido (to_image raises) → file still written (fallback path)
  - Handles empty findings list
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest


# ---------------------------------------------------------------------------
# Helpers — build the minimal Mission / schema objects the builders need
# ---------------------------------------------------------------------------
def _make_mission():
    """Return a lightweight mock Mission that satisfies the builders."""
    from types import SimpleNamespace
    return SimpleNamespace(
        id="DECK_TEST",
        site="ALL",
        user_query="deck smoke test",
        quest_id="quest:optimize_supply_chains",
        scope_tags=["fulfillment", "data_quality"],
        target_entity_kind="site",
        target_entity_key="ALL",
        horizon_days=90,
        status="done",
        progress_pct=100.0,
        artifact_paths={},
        parsed_intent={},
        created_at="2025-01-01T00:00:00",
        last_refreshed_at="2025-01-01T00:01:00",
        executed_by=None,
        owner_role="Planner",
    )


def _make_schema():
    from types import SimpleNamespace
    return SimpleNamespace(
        target_entity_kind="site",
        target_entity_key="ALL",
        tables=[],
        relationships=[],
        mermaid=None,
        as_dict=lambda: {"kind": "site", "tables": [], "relationships": []},
    )


@pytest.mark.unit
@pytest.mark.quest
class TestOnePager:
    def test_creates_file(self, tmp_path, synth_result):
        from src.brain.viz_composer import compose
        from src.deck.one_pager import render_one_pager
        mission = _make_mission()
        schema = _make_schema()
        viz = compose(synth_result)
        out = tmp_path / "one_pager.pptx"
        render_one_pager(mission, synth_result, viz, schema, out)
        assert out.exists(), "one_pager.pptx not created"

    def test_file_is_valid_pptx(self, tmp_path, synth_result):
        from pptx import Presentation
        from src.brain.viz_composer import compose
        from src.deck.one_pager import render_one_pager
        mission = _make_mission()
        schema = _make_schema()
        viz = compose(synth_result)
        out = tmp_path / "one_pager.pptx"
        render_one_pager(mission, synth_result, viz, schema, out)
        prs = Presentation(str(out))
        assert len(prs.slides) >= 1

    def test_survives_missing_kaleido(self, tmp_path, synth_result, monkeypatch):
        """Stub to_image to raise — fallback text path must still produce a file."""
        import plotly.basedatatypes
        monkeypatch.setattr(
            plotly.basedatatypes.BaseFigure, "to_image",
            lambda *a, **k: (_ for _ in ()).throw(
                RuntimeError("kaleido not available")),
        )
        from src.brain.viz_composer import compose
        from src.deck.one_pager import render_one_pager
        mission = _make_mission()
        schema = _make_schema()
        viz = compose(synth_result)
        out = tmp_path / "one_pager_nokaleido.pptx"
        render_one_pager(mission, synth_result, viz, schema, out)
        assert out.exists()

    def test_empty_findings(self, tmp_path):
        """Deck must not crash on a result with zero findings."""
        from src.brain.orchestrator import MissionResult
        from src.brain.viz_composer import compose
        from src.deck.one_pager import render_one_pager
        empty = MissionResult(
            mission_id="EMPTY", site="X", scope_tags=[],
            outcomes=[], findings=[],
            kpi_snapshot={}, progress_pct=0.0, elapsed_ms=0, refresh=False,
        )
        viz = compose(empty)
        mission = _make_mission()
        schema = _make_schema()
        out = tmp_path / "one_pager_empty.pptx"
        render_one_pager(mission, empty, viz, schema, out)
        assert out.exists()


@pytest.mark.unit
@pytest.mark.quest
class TestImplementationPlan:
    def test_creates_file(self, tmp_path, synth_result):
        from src.brain.viz_composer import compose
        from src.deck.implementation_plan import render_implementation_plan
        mission = _make_mission()
        schema = _make_schema()
        viz = compose(synth_result)
        out = tmp_path / "impl_plan.pptx"
        render_implementation_plan(mission, synth_result, viz, schema, out)
        assert out.exists(), "impl_plan.pptx not created"

    def test_file_is_valid_pptx(self, tmp_path, synth_result):
        from pptx import Presentation
        from src.brain.viz_composer import compose
        from src.deck.implementation_plan import render_implementation_plan
        mission = _make_mission()
        schema = _make_schema()
        viz = compose(synth_result)
        out = tmp_path / "impl_plan.pptx"
        render_implementation_plan(mission, synth_result, viz, schema, out)
        prs = Presentation(str(out))
        assert len(prs.slides) >= 1

    def test_survives_missing_kaleido(self, tmp_path, synth_result, monkeypatch):
        import plotly.basedatatypes
        monkeypatch.setattr(
            plotly.basedatatypes.BaseFigure, "to_image",
            lambda *a, **k: (_ for _ in ()).throw(
                RuntimeError("kaleido not available")),
        )
        from src.brain.viz_composer import compose
        from src.deck.implementation_plan import render_implementation_plan
        mission = _make_mission()
        schema = _make_schema()
        viz = compose(synth_result)
        out = tmp_path / "impl_plan_nokaleido.pptx"
        render_implementation_plan(mission, synth_result, viz, schema, out)
        assert out.exists()

    def test_empty_findings(self, tmp_path):
        from src.brain.orchestrator import MissionResult
        from src.brain.viz_composer import compose
        from src.deck.implementation_plan import render_implementation_plan
        empty = MissionResult(
            mission_id="EMPTY2", site="X", scope_tags=[],
            outcomes=[], findings=[],
            kpi_snapshot={}, progress_pct=0.0, elapsed_ms=0, refresh=False,
        )
        viz = compose(empty)
        mission = _make_mission()
        schema = _make_schema()
        out = tmp_path / "impl_plan_empty.pptx"
        render_implementation_plan(mission, empty, viz, schema, out)
        assert out.exists()
