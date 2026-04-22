"""
PPTX renderer (spec §8b).

render_pptx(findings, out_path) — produces the 16-slide portfolio deck or
14-slide single-site deck based on scope. Slide order is fixed and
deliberate (see spec §8b). This renderer reads the findings JSON produced
by findings.build_findings — it contains NO analytical logic, per spec §10.

Requires python-pptx. A clear ImportError is raised at entry if missing.
"""
from __future__ import annotations
import json
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    _PPTX_OK = True
except ImportError:  # pragma: no cover
    _PPTX_OK = False
    Presentation = Inches = Pt = Emu = RGBColor = MSO_SHAPE = PP_ALIGN = None  # type: ignore

SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5


def _require_pptx():
    if not _PPTX_OK:
        raise ImportError(
            "python-pptx is not installed. Install with:\n"
            "    pip install python-pptx>=0.6.23\n"
            "Then re-run the deck command."
        )


# Deck palette — instantiated lazily so this module is importable even when
# python-pptx is missing (the caller only hits _require_pptx at render time).
def _palette():
    return {
        "BG":     RGBColor(0xFF, 0xFF, 0xFF),
        "INK":    RGBColor(0x1A, 0x1A, 0x1A),
        "MUTED":  RGBColor(0x6B, 0x72, 0x80),
        "ACCENT": RGBColor(0x16, 0x4E, 0x87),
        "GOOD":   RGBColor(0x1F, 0x76, 0x3C),
        "WARN":   RGBColor(0xB3, 0x5C, 0x00),
        "BAD":    RGBColor(0xA6, 0x1B, 0x1B),
    }


# Placeholders reassigned inside render_pptx() before any builder runs.
C_BG = C_INK = C_MUTED = C_ACCENT = C_GOOD = C_WARN = C_BAD = None


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _new_prs(template_path: str = None) -> "Presentation":
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width  = Inches(SLIDE_W_IN)
        prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def _blank(prs):
    layout = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(layout)


def _box(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or C_BG
    shape.line.color.rgb = line or C_BG
    shape.line.width = Emu(0) if line is None else Pt(0.75)
    return shape


def _text(slide, left, top, width, height, text,
          size=14, bold=False, color=None, align=None):
    color = color or C_INK
    align = align if align is not None else PP_ALIGN.LEFT
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def _bullets(slide, left, top, width, height, items: list[str], size=12, color=None):
    color = color or C_INK
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def _hrule(slide, left, top, width, color=None, weight=1.5):
    color = color or C_ACCENT
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(weight))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.color.rgb = color
    return line


def _header(slide, title: str, eyebrow: str = "Cross-Dataset Supply-Chain Review"):
    _text(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.3),
          eyebrow, size=10, color=C_MUTED)
    _text(slide, Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.6),
          title, size=24, bold=True, color=C_INK)
    _hrule(slide, Inches(0.5), Inches(1.18), Inches(12.3))


def _kpi_card(slide, left, top, label: str, kpi: dict, width=Inches(3.9), height=Inches(2.0)):
    _box(slide, left, top, width, height, fill=C_BG, line=C_MUTED)
    _text(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.35),
          label.upper(), size=10, bold=True, color=C_MUTED)
    value = kpi.get("value")
    shown = "—" if value is None else f"{value}%"
    _text(slide, left + Inches(0.2), top + Inches(0.45), width - Inches(0.4), Inches(0.7),
          shown, size=36, bold=True, color=C_INK)
    delta = kpi.get("delta_pp")
    if delta is None:
        dtxt = "Δ vs prior 14d: n/a"
        dc = C_MUTED
    else:
        sign = "+" if delta >= 0 else ""
        dtxt = f"Δ vs prior 14d: {sign}{delta} pp"
        dc = C_GOOD if delta >= 0 else C_BAD
    _text(slide, left + Inches(0.2), top + Inches(1.15), width - Inches(0.4), Inches(0.3),
          dtxt, size=11, bold=True, color=dc)
    goal = kpi.get("goal"); n = kpi.get("n"); base = kpi.get("baseline_90d")
    meta = f"goal {goal}%   |   90d {base if base is not None else '—'}%   |   n={n:,}"
    _text(slide, left + Inches(0.2), top + Inches(1.5), width - Inches(0.4), Inches(0.3),
          meta, size=9, color=C_MUTED)


def _table(slide, left, top, width, height, headers: list[str], rows: list[list]):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(h)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(10); run.font.bold = True
                run.font.color.rgb = C_BG
        cell.fill.solid(); cell.fill.fore_color.rgb = C_ACCENT
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = "" if val is None else str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(10); run.font.color.rgb = C_INK
    return table


# ---------------------------------------------------------------------------
# Slide builders (fixed order per §8b)
# ---------------------------------------------------------------------------

def _slide_01_cover(prs, f):
    s = _blank(prs)
    _box(s, 0, 0, prs.slide_width, prs.slide_height, fill=C_BG)
    _box(s, 0, Inches(6.7), prs.slide_width, Inches(0.8), fill=C_ACCENT)
    _text(s, Inches(0.5), Inches(1.5), Inches(12), Inches(0.5),
          "Cross-Dataset Supply-Chain Review", size=14, color=C_MUTED)
    scope = f["scope"]
    title = "Portfolio View" if scope["site"] == "ALL" else f"Site: {scope['site']}"
    _text(s, Inches(0.5), Inches(2.1), Inches(12), Inches(1.3),
          title, size=44, bold=True, color=C_INK)
    win = scope["windows"]
    sub = (f"Anchor T = {win['T']}   |   "
           f"Past 14d {win['past_14d'][0]} → {win['past_14d'][1]}   |   "
           f"ERP: {scope.get('erp','mixed')}")
    _text(s, Inches(0.5), Inches(3.6), Inches(12), Inches(0.5),
          sub, size=14, color=C_MUTED)
    datasets = "Datasets in scope: " + " · ".join(scope["datasets"])
    _text(s, Inches(0.5), Inches(4.1), Inches(12), Inches(0.4),
          datasets, size=12, color=C_MUTED)
    _text(s, Inches(0.5), Inches(6.85), Inches(12), Inches(0.6),
          f"Anchor policy {scope['anchor_policy']}   ·   Reproducibility seed {scope['seed']}",
          size=11, color=C_BG)


def _slide_02_exec(prs, f):
    s = _blank(prs)
    _header(s, "Executive summary")
    # Three KPI cards
    _kpi_card(s, Inches(0.5), Inches(1.5), "OTD 14d", f["kpis"]["otd"])
    _kpi_card(s, Inches(4.7), Inches(1.5), "IFR 14d", f["kpis"]["ifr"])
    _kpi_card(s, Inches(8.9), Inches(1.5), "Cycle-count 14d", f["kpis"]["cc"])
    # Four lenses
    _text(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.3),
          "FOUR LENSES", size=10, bold=True, color=C_MUTED)
    lenses = [
        f"Failure signature — top OTD reason: {_top_reason(f)}",
        f"Allocation vs. stockout — alloc-gap {f['failure_signatures']['ifr'].get('allocation_gap_pct','—')}% "
        f"| hard-stockout {f['failure_signatures']['ifr'].get('hard_stockout_pct','—')}%",
        f"PFEP match rate on miss parts — {_pct(f['intersections']['pfep_match'].get('match_rate'))}",
        f"Triple intersection (CC∩IFR∩OTD, 90d) — {f['intersections']['triple']['triple_count']} parts",
    ]
    _bullets(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(1.6), lenses, size=12)
    # Top 3 realizations
    _text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.3),
          "REALIZATIONS", size=10, bold=True, color=C_MUTED)
    rl = [r["realization"] for r in f["realizations"][:3]] or ["(no rules fired — data may be too sparse)"]
    _bullets(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.3), rl, size=12, color=C_ACCENT)


def _slide_03_architecture(prs, f):
    s = _blank(prs)
    _header(s, "Data architecture — four datasets joined through Part No.")
    labels = [("OTD", "sales-order line history"),
              ("IFR", "order-line snapshot at order time"),
              ("ITR", "inventory transactions (cycle-count)"),
              ("PFEP", "plan-for-every-part master")]
    # 2x2 grid
    coords = [(0.9, 1.7), (7.4, 1.7), (0.9, 4.4), (7.4, 4.4)]
    for (ds, sub), (x, y) in zip(labels, coords):
        _box(s, Inches(x), Inches(y), Inches(5.0), Inches(2.4), fill=C_BG, line=C_ACCENT)
        _text(s, Inches(x + 0.2), Inches(y + 0.15), Inches(4.6), Inches(0.5),
              ds, size=22, bold=True, color=C_ACCENT)
        _text(s, Inches(x + 0.2), Inches(y + 0.75), Inches(4.6), Inches(0.5),
              sub, size=12, color=C_MUTED)
        fields = _top_fields(ds)
        _bullets(s, Inches(x + 0.2), Inches(y + 1.15), Inches(4.6), Inches(1.2),
                 fields, size=10, color=C_INK)
    # Central PART key callout
    _box(s, Inches(6.1), Inches(3.4), Inches(1.1), Inches(1.1), fill=C_ACCENT)
    _text(s, Inches(6.1), Inches(3.75), Inches(1.1), Inches(0.4),
          "PART", size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)


def _top_fields(ds: str) -> list[str]:
    from .schemas import DATASETS
    return DATASETS[ds]["required"][:5]


def _slide_04_otd_scorecard(prs, f):
    s = _blank(prs)
    _header(s, "OTD scorecard")
    _kpi_card(s, Inches(0.5), Inches(1.5), "OTD 14d", f["kpis"]["otd"], width=Inches(5.0), height=Inches(2.4))
    tail = f["failure_signatures"]["days_late"]
    items = [
        f"Median Days Late: {tail.get('median') or '—'}",
        f"p75 / p90 / p99: {tail.get('p75')} / {tail.get('p90')} / {tail.get('p99')}",
        f"Max Days Late: {tail.get('max') or '—'}",
        f"Fat-tail flag (p99 > 30): {'YES' if tail.get('fat_tail') else 'no'}",
    ]
    _text(s, Inches(6.0), Inches(1.5), Inches(7.0), Inches(0.4),
          "DAYS-LATE DISTRIBUTION", size=10, bold=True, color=C_MUTED)
    _bullets(s, Inches(6.0), Inches(1.85), Inches(7.0), Inches(2.2), items, size=12)

    by_site = f.get("by_site") or {}
    if by_site:
        rows = []
        for site, kpi in list(by_site.items())[:10]:
            if not isinstance(kpi, dict) or "otd" not in kpi:
                continue
            rows.append([site, _pct(kpi["otd"].get("value"), 1),
                         _pp(kpi["otd"].get("delta_pp")),
                         _pct(kpi["otd"].get("baseline_90d"), 1),
                         kpi["otd"].get("n")])
        if rows:
            _text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.3),
                  "PER-SITE OTD", size=10, bold=True, color=C_MUTED)
            _table(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
                   ["Site", "OTD 14d", "Δ pp", "90d", "n"], rows)


def _slide_05_otd_signature(prs, f):
    s = _blank(prs)
    _header(s, "OTD failure signature")
    sig = f["failure_signatures"]["otd"]
    _text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4),
          f"Total late lines (14d): {sig['total_late_lines']}", size=13, color=C_INK)
    rows = [[r["reason"], r["count"], f"{r['pct']}%", r["class"]] for r in sig["top_reasons"]]
    if rows:
        _table(s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(4.0),
               ["Failure reason", "Count", "Share", "Class"], rows)
    _text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
          "Class = EXECUTION / KITTING / SCHEDULING / SUPPLY — routes to pathway templates.",
          size=10, color=C_MUTED)


def _slide_06_otd_customers(prs, f):
    s = _blank(prs)
    _header(s, "OTD — customer centrality + days-late tail")
    custs = f["centrality"]["customers"]
    rows = [[c["customer"], c["late_lines"], f"{c['pct']}%",
             "★" if c["concentrated"] else ""] for c in custs]
    if rows:
        _table(s, Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.0),
               ["Customer", "Late lines", "Share", "Flag"], rows)
    tail = f["failure_signatures"]["days_late"]
    items = [
        f"n = {tail.get('n',0)} late lines",
        f"median = {tail.get('median') or '—'}",
        f"p75 = {tail.get('p75') or '—'}",
        f"p90 = {tail.get('p90') or '—'}",
        f"p99 = {tail.get('p99') or '—'}",
        f"max = {tail.get('max') or '—'}",
    ]
    _text(s, Inches(8.5), Inches(1.5), Inches(4.5), Inches(0.4),
          "DAYS-LATE TAIL", size=10, bold=True, color=C_MUTED)
    _bullets(s, Inches(8.5), Inches(1.9), Inches(4.5), Inches(3.0), items, size=12)


def _slide_07_ifr_decomp(prs, f):
    s = _blank(prs)
    _header(s, "IFR scorecard + miss decomposition")
    _kpi_card(s, Inches(0.5), Inches(1.5), "IFR 14d", f["kpis"]["ifr"],
              width=Inches(5.0), height=Inches(2.4))
    d = f["failure_signatures"]["ifr"]
    # Stacked decomposition bar, drawn as three colored boxes.
    x0 = Inches(6.0); y0 = Inches(1.8); w = Inches(6.8); h = Inches(0.8)
    total = d["total_miss"] or 1
    parts = [
        ("allocation_gap", d["allocation_gap"], C_ACCENT, "Allocation gap"),
        ("hard_stockout",  d["hard_stockout"],  C_BAD,    "Hard stockout"),
        ("covered_miss",   d["covered_miss"],   C_MUTED,  "Covered miss"),
    ]
    cursor = x0
    for _, n, color, label in parts:
        width = Emu(int(w * n / total))
        if width <= 0:
            continue
        bar = _box(s, cursor, y0, width, h, fill=color, line=color)
        cursor += width
    _text(s, Inches(6.0), Inches(2.7), Inches(6.8), Inches(0.3),
          f"n={d['total_miss']} misses   |   alloc-gap {d['allocation_gap_pct']}%   "
          f"hard-stockout {d['hard_stockout_pct']}%   covered {d['covered_miss_pct']}%",
          size=10, color=C_MUTED)
    _text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.4),
          "DIAGNOSTIC INTERPRETATION", size=10, bold=True, color=C_MUTED)
    interp = [
        "High allocation-gap → reservation/pegging issue (Oracle hygiene).",
        "High hard-stockout → safety-stock problem (PFEP fix).",
        "Covered-miss should be near zero — if not, investigate data quality.",
    ]
    _bullets(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.0), interp, size=12)


def _slide_08_alloc_thesis(prs, f):
    s = _blank(prs)
    _header(s, "IFR allocation-gap thesis")
    d = f["failure_signatures"]["ifr"]
    pct = d.get("allocation_gap_pct") or 0
    color = C_BAD if pct > 30 else (C_WARN if pct > 20 else C_GOOD)
    _text(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.2),
          f"{pct}% of IFR misses are allocation-gap, not stockout.",
          size=28, bold=True, color=color)
    _text(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.4),
          "MECHANISM", size=10, bold=True, color=C_MUTED)
    _bullets(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(3.0), [
        "On Hand Qty > 0 AND Available Qty < SO Qty → stock is there but reserved elsewhere.",
        "Root cause is reservation logic (pegging, MRP allocations, stale orders), not supply.",
        "Fixing via safety-stock raise is a mis-diagnosis — inventory grows without OTIF lift.",
        "Systemic pathway #4 — nightly allocation reconciliation batch job.",
    ], size=13)


def _slide_09_vendors(prs, f):
    s = _blank(prs)
    _header(s, "Vendor centrality — shared leverage on OTD + IFR")
    rows = [[v["supplier"], v["combined"], v["otd_late"], v["ifr_miss"]]
            for v in f["centrality"]["vendors"]]
    if rows:
        _table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.0),
               ["Supplier", "Combined (late+miss)", "OTD late", "IFR miss"], rows)
    _text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8),
          "Top-5 is the intervention list — focused scorecard, not full supplier base.",
          size=11, color=C_MUTED)


def _slide_10_cc(prs, f):
    s = _blank(prs)
    _header(s, "Cycle-count cadence + variance")
    cc = f["failure_signatures"]["cc"]
    _kpi_card(s, Inches(0.5), Inches(1.5), "CC accuracy 14d", f["kpis"]["cc"],
              width=Inches(5.0), height=Inches(2.4))
    items = [
        f"Active days / business days: {cc['active_days']} / {cc['business_days']}  "
        f"({cc['cadence_compliance_pct']}%)",
        f"Absolute variance $: ${cc['absolute_variance_$']:,}",
        f"Net variance $: ${cc['net_variance_$']:,}  (positive = found-stock bias)",
        f"Repeat-offender share: {cc['repeat_offender_pct']}% of tx",
        f"Reason Code populated: {cc['reason_code_populated_pct']}%",
    ]
    _text(s, Inches(6.0), Inches(1.5), Inches(7.0), Inches(0.4),
          "CADENCE & VARIANCE", size=10, bold=True, color=C_MUTED)
    _bullets(s, Inches(6.0), Inches(1.85), Inches(7.0), Inches(3.0), items, size=12)


def _slide_11_pfep(prs, f):
    s = _blank(prs)
    _header(s, "PFEP parameter health")
    audit = f["intersections"]["pfep_match"]
    rec = f["intersections"]["recoverable"]
    _text(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4),
          f"Miss parts matched against PFEP: {_pct(audit.get('match_rate'))}  "
          f"(n={audit.get('miss_n')})", size=14, bold=True, color=C_INK)
    bullets = [
        f"Safety Stock missing on matched: {audit.get('ss_missing_pct')}%",
        f"ABC missing on matched: {audit.get('abc_missing_pct')}%",
        f"Lead Time = 0 on matched: {audit.get('lt_zero_pct')}%",
    ]
    _bullets(s, Inches(0.5), Inches(2.1), Inches(12.3), Inches(1.5), bullets, size=12)
    if rec.get("blocked_by_match_rate"):
        _text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.5),
              f"RECOVERABLE STOCKOUT CALC BLOCKED: match rate {_pct(rec['match_rate'])} < 90%.",
              size=14, bold=True, color=C_BAD)
        _text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(1.5),
              "Resolve canonical part-number alignment before acting on recoverable count.",
              size=12, color=C_MUTED)
    else:
        _text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.4),
              "RECOVERABLE STOCKOUTS", size=10, bold=True, color=C_MUTED)
        _text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.8),
              f"{rec.get('recoverable')} of {rec.get('purchased_stockouts')} purchased "
              f"stockouts are PFEP-preventable ({rec.get('recoverable_pct')}%).",
              size=16, bold=True, color=C_ACCENT)


def _slide_12_intersections(prs, f):
    s = _blank(prs)
    _header(s, "Cross-dataset intersection — where fixes compound")
    t = f["intersections"]["triple"]
    rows = [
        ["CC only",            t["cc_only"]],
        ["IFR only",           t["ifr_only"]],
        ["OTD only",           t["otd_only"]],
        ["CC ∩ IFR",           t["cc_and_ifr"]],
        ["CC ∩ OTD",           t["cc_and_otd"]],
        ["OTD ∩ IFR",          t["otd_and_ifr"]],
        ["CC ∩ IFR ∩ OTD",     t["triple_count"]],
    ]
    _table(s, Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.5),
           ["Bucket", "Parts (90d)"], rows)
    if t["triple_count"]:
        _text(s, Inches(6.5), Inches(1.5), Inches(6.3), Inches(0.4),
              "TRIPLE-INTERSECTION PARTS (first 15)", size=10, bold=True, color=C_MUTED)
        _bullets(s, Inches(6.5), Inches(1.85), Inches(6.3), Inches(5.0),
                 t["triple"][:15], size=11)
    else:
        _text(s, Inches(6.5), Inches(1.5), Inches(6.3), Inches(4.0),
              "No parts appear in all three datasets.\nDatasets pick up different failure "
              "modes — interventions run in parallel.", size=13, color=C_MUTED)


def _slide_13_systemic(prs, f):
    s = _blank(prs)
    _header(s, "Systemic pathways — ERP-level configuration")
    rows = [[p["rule_id"], p["name"], p["owner"], p["sequence"]]
            for p in f["pathways_systemic"]]
    if rows:
        _table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.0),
               ["Trigger", "Pathway", "Owner", "Seq"], rows)
    _text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.4),
          "EXPECTED DOWNSTREAM LIFT", size=10, bold=True, color=C_MUTED)
    lifts = [f"{p['name']} — {p['downstream_lift']}" for p in f["pathways_systemic"]]
    _bullets(s, Inches(0.5), Inches(5.05), Inches(12.3), Inches(2.1), lifts or
             ["(no systemic pathway fired — extend rule table in §6 or check data population)"],
             size=11)


def _slide_14_operational(prs, f):
    s = _blank(prs)
    _header(s, "Operational pathways — plant-floor action cards")
    ops = f["pathways_operational"]
    rows = [[p.get("trigger_reason") or p.get("rule_id") or "—",
             p["name"], p["owner"]] for p in ops]
    if rows:
        _table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.0),
               ["Trigger", "Action", "Owner"], rows)
    _text(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
          "Operational fixes are per-site — pair each action with a named Site Lead before the "
          "next review.", size=11, color=C_MUTED)


def _slide_15_roadmap(prs, f):
    s = _blank(prs)
    _header(s, "30 · 60 · 90 · 180-day roadmap")
    rm = f["roadmap"]
    y = 1.5
    for phase in ("T+30", "T+60", "T+90", "T+180"):
        _text(s, Inches(0.5), Inches(y), Inches(1.5), Inches(0.4),
              phase, size=14, bold=True, color=C_ACCENT)
        _bullets(s, Inches(2.2), Inches(y), Inches(10.5), Inches(1.3),
                 rm.get(phase, []), size=11)
        y += 1.4


def _slide_16_governance(prs, f):
    s = _blank(prs)
    _header(s, "Governance + closing")
    gov = f["governance"]
    items = [
        f"Cadence: {gov['cadence']}",
        f"Targets: OTD {gov['targets']['OTD']}%, IFR {gov['targets']['IFR']}%, "
        f"CC accuracy ≥ {gov['targets']['CC']}%",
        f"Review runbook: {gov['review_runbook']}",
    ]
    _bullets(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.0), items, size=13)
    _text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(1.0),
          "Systemic fixes make tomorrow's review easier. Operational fixes make today's review honest.",
          size=16, bold=True, color=C_ACCENT)
    scope = f["scope"]
    _text(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.3),
          f"{scope['anchor_policy']}   ·   seed {scope['seed']}   ·   "
          f"generated for {scope['site']}",
          size=9, color=C_MUTED)


# ---------------------------------------------------------------------------
# Small formatting helpers
# ---------------------------------------------------------------------------

def _pct(x, digits=1):
    if x is None:
        return "—"
    if isinstance(x, float) and x <= 1.0:
        return f"{round(100*x, digits)}%"
    return f"{round(x, digits)}%"


def _pp(x):
    if x is None:
        return "—"
    sign = "+" if x >= 0 else ""
    return f"{sign}{x} pp"


def _top_reason(f) -> str:
    reasons = f["failure_signatures"]["otd"]["top_reasons"]
    return reasons[0]["reason"] if reasons else "—"


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def render_pptx(findings: dict, out_path: str | Path, template_path: str | Path = None) -> Path:
    """Render the findings payload to a .pptx file. Returns the output path."""
    _require_pptx()
    global C_BG, C_INK, C_MUTED, C_ACCENT, C_GOOD, C_WARN, C_BAD
    p = _palette()
    C_BG, C_INK, C_MUTED = p["BG"], p["INK"], p["MUTED"]
    C_ACCENT, C_GOOD, C_WARN, C_BAD = p["ACCENT"], p["GOOD"], p["WARN"], p["BAD"]

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = _new_prs(template_path)
    single_site = findings["scope"]["site"] != "ALL"

    # Fixed slide order per §8b. Single-site deck drops slides 4 (per-site
    # scorecard table collapses to a single-site view, already covered by
    # slide 2) and 8 (allocation-gap thesis is only drawn when alloc-gap
    # is meaningful — kept even for single-site).
    builders = [
        _slide_01_cover,
        _slide_02_exec,
        _slide_03_architecture,
        _slide_04_otd_scorecard,
        _slide_05_otd_signature,
        _slide_06_otd_customers,
        _slide_07_ifr_decomp,
        _slide_08_alloc_thesis,
        _slide_09_vendors,
        _slide_10_cc,
        _slide_11_pfep,
        _slide_12_intersections,
        _slide_13_systemic,
        _slide_14_operational,
        _slide_15_roadmap,
        _slide_16_governance,
    ]
    if single_site:
        # §8b — 14-slide single-site: drop per-site table (04) and alloc-gap
        # thesis (08) when alloc-gap ≤ 20%.
        alloc_pct = findings["failure_signatures"]["ifr"].get("allocation_gap_pct") or 0
        drop = {_slide_04_otd_scorecard}
        if alloc_pct <= 20:
            drop.add(_slide_08_alloc_thesis)
        builders = [b for b in builders if b not in drop]

    for build in builders:
        build(prs, findings)

    prs.save(out)
    return out


def dump_findings_json(findings: dict, out_path: str | Path) -> Path:
    """Phase 8a — write the findings payload as pretty JSON."""
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(findings, indent=2, default=str), encoding="utf-8")
    return out
