"""
Bi-Weekly 1 Pager renderer — single-slide executive summary.

Encodes the Brain's understanding of the Bi-Weekly 1 Pager template
(config/templates/Bi-Weekly 1 Pager.pptx):
  - Slide size: 13.33" × 7.50" landscape (widescreen)
  - Layout: blank master, no placeholders
  - Design: dark-navy header bar, three KPI tiles, four-zone body

Zone map (all measurements in inches):
┌──────────────────────────────────────────────────┐  0.00
│ HEADER: site · date range · "Bi-Weekly Review"   │  0.60
├──────────────┬───────────────────────────────────┤
│  KPI TILE 1  │  KPI TILE 2  │  KPI TILE 3       │  1.85
├──────────────┼────────────────────────────────────┤
│ FOUR LENSES  │  REALIZATIONS                     │  3.80
│ (left half)  │  + TOP 3 ACTIONS  (right half)    │  6.45
├──────────────┴────────────────────────────────────┤
│ FOOTER: policy · seed · site · timestamp         │  7.50
└──────────────────────────────────────────────────┘

Public API:
    render_biweekly_one_pager(findings, out_path, template_path=None) -> Path
"""
from __future__ import annotations

import io
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

# ── python-pptx lazy import ────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False
    Presentation = Inches = Pt = Emu = RGBColor = PP_ALIGN = None  # type: ignore

SLIDE_W = 13.333
SLIDE_H = 7.5

# ── Brain colour palette (matches SCB dark-mode aesthetic) ────────────────
_C = {
    "bg":      (0xFF, 0xFF, 0xFF),      # slide background
    "ink":     (0x1A, 0x1A, 0x1A),      # body text
    "muted":   (0x6B, 0x72, 0x80),      # labels / captions
    "navy":    (0x0C, 0x1E, 0x3B),      # header bar fill
    "accent":  (0x16, 0x4E, 0x87),      # accent blue
    "good":    (0x1F, 0x76, 0x3C),      # green KPI
    "warn":    (0xB3, 0x5C, 0x00),      # amber KPI
    "bad":     (0xA6, 0x1B, 0x1B),      # red KPI
    "tile_bg": (0xF1, 0xF5, 0xF9),      # tile fill
    "rule":    (0xCB, 0xD5, 0xE1),      # horizontal rules
    "insight": (0xEF, 0xF6, 0xFF),      # insight panel fill
}

def _rgb(*key_or_tuple):
    if len(key_or_tuple) == 1:
        v = _C[key_or_tuple[0]]
    else:
        v = key_or_tuple
    return RGBColor(*v)


# ── Low-level drawing helpers ───────────────────────────────────────────────

def _box(slide, l, t, w, h, fill="bg", line_color=None, line_pt=0.5):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(fill)
    if line_color:
        s.line.color.rgb = _rgb(line_color)
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s


def _text(slide, l, t, w, h, txt, size=11, bold=False, italic=False,
          color="ink", align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(txt)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    return txBox


def _bullet_block(slide, l, t, w, h, items: list[str],
                  size=10, label_size=9, label="", label_color="muted",
                  bullet_color="ink", fill=None):
    """Render a labelled bullet block, optionally with a background box."""
    if fill:
        _box(slide, l, t, w, h, fill=fill)
    if label:
        _text(slide, l + 0.08, t + 0.06, w - 0.16, 0.22,
              label.upper(), size=label_size, bold=True, color=label_color)
        t_start = t + 0.28
    else:
        t_start = t + 0.06
    txBox = slide.shapes.add_textbox(Inches(l + 0.08), Inches(t_start),
                                     Inches(w - 0.16), Inches(h - 0.28))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = _rgb(bullet_color)


def _kpi_tile(slide, l, t, w, h, label, value, delta, sub, status="muted"):
    """Render a KPI tile: label / large value / delta / sub-text."""
    _box(slide, l, t, w, h, fill="tile_bg", line_color="rule")
    # label
    _text(slide, l + 0.12, t + 0.10, w - 0.24, 0.22,
          label, size=9, bold=True, color="muted")
    # value
    val_color = {"good": "good", "warn": "warn", "bad": "bad"}.get(status, "ink")
    _text(slide, l + 0.12, t + 0.32, w - 0.24, 0.55,
          str(value), size=26, bold=True, color=val_color)
    # delta
    _text(slide, l + 0.12, t + 0.88, w - 0.24, 0.22,
          str(delta), size=9, color="muted")
    # sub (goal / 90d / n)
    _text(slide, l + 0.12, t + 1.10, w - 0.24, 0.22,
          str(sub), size=8, color="muted")


# ── Status helper ────────────────────────────────────────────────────────────

def _kpi_status(value, goal=95.0):
    if value is None or (isinstance(value, float) and (value != value)):  # nan
        return "muted"
    if value >= goal:
        return "good"
    if value >= goal * 0.90:
        return "warn"
    return "bad"


def _fmt(v, suffix="%", prec=1):
    if v is None or (isinstance(v, float) and (v != v)):
        return "—"
    return f"{v:.{prec}f}{suffix}"


def _delta_str(now, prior):
    if now is None or prior is None:
        return "Δ vs prior: n/a"
    if (isinstance(now, float) and now != now) or (isinstance(prior, float) and prior != prior):
        return "Δ vs prior: n/a"
    d = now - prior
    sign = "+" if d >= 0 else ""
    return f"Δ vs prior 14d: {sign}{d:.1f} pp"


# ── Brain insight generator ──────────────────────────────────────────────────

def _generate_insights(f: dict) -> list[str]:
    """
    Generate up to 6 Brain-derived insights from the findings dict.
    Uses the actual findings schema:
      f["kpis"]["otd|ifr|cc"]["value|prior|delta_pp|baseline_90d|n"]
      f["failure_signatures"]["ifr"]["hard_stockout_pct|allocation_gap_pct"]
      f["failure_signatures"]["cc"]["active_days|business_days|reason_code_populated_pct"]
      f["intersections"]["pfep_match"]["ss_missing_pct|abc_missing_pct|lt_zero_pct"]
      f["intersections"]["recoverable"]["purchased_stockouts|recoverable"]
      f["intersections"]["triple"]["triple_count"]
      f["centrality"]["vendors"][0]["supplier|otd_late|combined"]
    """
    insights = []

    kpis  = f.get("kpis", {})
    fsig  = f.get("failure_signatures", {})
    inter = f.get("intersections", {})
    cent  = f.get("centrality", {})

    ifr_kpi = kpis.get("ifr", {})
    ifr_14  = ifr_kpi.get("value")
    if ifr_14 is not None and ifr_14 == ifr_14 and ifr_14 < 95:
        gap = 95 - ifr_14
        ifr_sig = fsig.get("ifr", {})
        stockout_pct = ifr_sig.get("hard_stockout_pct", 0) or 0
        alloc_pct    = ifr_sig.get("allocation_gap_pct", 0) or 0
        if stockout_pct > 50:
            insights.append(
                f"IFR {ifr_14:.1f}% is {gap:.1f} pp below goal — "
                f"{stockout_pct:.0f}% of misses are hard stockouts. "
                "Safety-stock population in Oracle PFEP closes this structurally."
            )
        elif alloc_pct > 30:
            insights.append(
                f"IFR {ifr_14:.1f}% gap ({gap:.1f} pp) is allocation-driven — "
                f"{alloc_pct:.0f}% of misses are pegging/reservation failures. "
                "Oracle allocation hygiene batch resolves."
            )
        else:
            insights.append(
                f"IFR {ifr_14:.1f}% is {gap:.1f} pp below the 95% goal. "
                f"Alloc gap: {alloc_pct:.0f}%  |  Hard stockout: {stockout_pct:.0f}%."
            )

    # OTD trend vs 90d baseline
    otd_kpi = kpis.get("otd", {})
    otd_14  = otd_kpi.get("value")
    otd_90  = otd_kpi.get("baseline_90d")
    if otd_14 is not None and otd_90 is not None:
        if otd_14 == otd_14 and otd_90 == otd_90:
            trend = otd_14 - otd_90
            if abs(trend) >= 2:
                dir_word = "improving" if trend > 0 else "deteriorating"
                insights.append(
                    f"OTD is {dir_word} vs 90-day baseline "
                    f"({otd_90:.1f}% → {otd_14:.1f}%, "
                    f"{'+'if trend>0 else ''}{trend:.1f} pp)."
                )

    # Recoverable stockouts via PFEP
    rec = inter.get("recoverable", {})
    purchased = rec.get("purchased_stockouts")
    recoverable = rec.get("recoverable")
    if recoverable is not None and purchased is not None and recoverable > 0:
        pct = 100 * recoverable / max(purchased, 1)
        insights.append(
            f"{recoverable:,} of {purchased:,} stockouts ({pct:.0f}%) are "
            "PFEP-preventable — one-time data population = compounding fill-rate ROI."
        )
    else:
        pm = inter.get("pfep_match", {})
        match_rate = pm.get("match_rate", 0) or 0
        ss_missing = pm.get("ss_missing_pct", 0) or 0
        if match_rate >= 0.8 and ss_missing > 50:
            insights.append(
                f"PFEP match rate {match_rate*100:.0f}% on miss parts — "
                f"Safety Stock missing on {ss_missing:.0f}%. "
                "Filling SS triggers formula-driven replenishment immediately."
            )

    # CC cadence
    cc_sig   = fsig.get("cc", {})
    cc_days  = cc_sig.get("active_days", 0) or 0
    biz_days = cc_sig.get("business_days", 10) or 10
    if cc_days == 0:
        insights.append(
            f"Cycle-count program stopped "
            f"(0 of {biz_days} business days active). "
            "All inventory accuracy KPIs are directional only until cadence restores."
        )
    elif cc_days < biz_days * 0.6:
        insights.append(
            f"Cycle-count cadence {cc_days}/{biz_days} days ({100*cc_days/biz_days:.0f}%) — "
            "below 60% compliance threshold. Variance attribution unreliable."
        )

    # Reason-code completeness
    rcode = cc_sig.get("reason_code_populated_pct", None)
    if rcode is not None and rcode < 50:
        insights.append(
            f"Variance reason codes populated on only {rcode:.0f}% of adjustments. "
            "Root-cause targeting is impossible until field enforcement is added."
        )

    # Vendor concentration
    vendors = cent.get("vendors", []) or []
    if vendors:
        top = vendors[0]
        top_name  = top.get("supplier", "—")
        top_comb  = top.get("combined", 0) or 0
        top_late  = top.get("otd_late", 0) or 0
        if top_comb >= 3:
            insights.append(
                f"Top OTD+IFR leverage vendor: {top_name} "
                f"({top_late} late lines, {top_comb} combined misses). "
                "Focused scorecard session = fastest recovery path."
            )

    # Triple intersection
    triple_count = inter.get("triple", {}).get("triple_count", 0) or 0
    if triple_count > 0:
        insights.append(
            f"{triple_count} part(s) appear in CC∩IFR∩OTD failure sets — "
            "highest-leverage intervention targets."
        )

    return insights[:6]


# ── 30-day action encoding ────────────────────────────────────────────────────

def _top_actions(f: dict) -> list[str]:
    """Return the top 3 most urgent 30-day actions derived from findings."""
    actions = []
    pm = f.get("intersections", {}).get("pfep_match", {})
    ss_missing = pm.get("ss_missing_pct", 0) or 0
    abc_missing = pm.get("abc_missing_pct", 0) or 0
    lt_zero = pm.get("lt_zero_pct", 0) or 0

    if ss_missing > 50:
        actions.append(
            f"Close PFEP gaps: Safety Stock missing on {ss_missing:.0f}% of miss parts — "
            "populate to unlock formula-driven replenishment."
        )
    if abc_missing > 50:
        actions.append(
            f"Populate ABC classification ({abc_missing:.0f}% missing) — "
            "required by AST-INV-PRO-0001 to drive cycle-count cadence."
        )
    if lt_zero > 50:
        actions.append(
            f"Fix Lead Time = 0 on {lt_zero:.0f}% of parts — "
            "zero LT causes infinite MRP urgency; replaces demand signal noise."
        )

    pathways = f.get("pathways_systemic", []) or []
    seen = set()
    for pw in pathways:
        name = pw.get("name", "")
        if name and name not in seen:
            actions.append(f"{name} — {pw.get('downstream_lift', '')[:120]}")
            seen.add(name)
        if len(actions) >= 5:
            break

    # Fallback if findings are sparse
    if not actions:
        actions = [
            "Populate ABC Inventory Catalog — drives CC cadence per AST-INV-PRO-0001.",
            "Switch Safety Stock Planning Method — converts recoverable stockouts to fills.",
            "Enforce Transaction Reason Code on cycle-count adjustments.",
        ]
    return [a for a in actions if a][:3]


# ── Main renderer ─────────────────────────────────────────────────────────────

def render_biweekly_one_pager(
    findings: dict,
    out_path: str | Path,
    template_path: str | Path | None = None,
) -> Path:
    """
    Render the Bi-Weekly 1 Pager for the given findings dict.

    Parameters
    ----------
    findings     : output of src.deck.findings.build_findings()
    out_path     : destination .pptx file path
    template_path: optional path to Bi-Weekly 1 Pager.pptx template
                   (uses template master/theme when provided)

    Returns
    -------
    Path to the written file.
    """
    if not _PPTX_OK:
        raise ImportError("python-pptx is required: pip install python-pptx>=0.6.23")

    # ── Build presentation ────────────────────────────────────────────────────
    # Always start from a fresh Presentation() to avoid the python-pptx
    # duplicate-slide-name ZIP corruption that occurs when opening a template
    # that already has a slide1.xml and then calling add_slide() (which also
    # writes slide1.xml).  We read the template only to extract the widescreen
    # dimensions so the output matches the corporate master exactly.
    tmpl = Path(template_path) if template_path else None
    if tmpl and tmpl.exists():
        try:
            _ref = Presentation(str(tmpl))
            _w   = _ref.slide_width
            _h   = _ref.slide_height
        except Exception:
            _w = Inches(SLIDE_W)
            _h = Inches(SLIDE_H)
    else:
        _w = Inches(SLIDE_W)
        _h = Inches(SLIDE_H)

    prs = Presentation()
    prs.slide_width  = _w
    prs.slide_height = _h
    # Layout 6 is the blank layout in every built-in theme
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Extract findings values ───────────────────────────────────────────────
    scope   = findings.get("scope", {})
    kpis    = findings.get("kpis", {})
    site    = scope.get("site", "—") or "—"
    windows = scope.get("windows", {})
    past_14 = windows.get("past_14d", [])
    win_lo  = past_14[0] if len(past_14) > 0 else ""
    win_hi  = past_14[1] if len(past_14) > 1 else ""
    anchor  = windows.get("T", "")

    def _kv(kpi_name, key, default=None):
        return kpis.get(kpi_name, {}).get(key, default)

    otd_14  = _kv("otd", "value")
    otd_90  = _kv("otd", "baseline_90d")
    otd_n   = _kv("otd", "n", 0) or 0
    otd_p   = _kv("otd", "prior")
    ifr_14  = _kv("ifr", "value")
    ifr_90  = _kv("ifr", "baseline_90d")
    ifr_n   = _kv("ifr", "n", 0) or 0
    ifr_p   = _kv("ifr", "prior")
    cc_14   = _kv("cc", "value")
    cc_90   = _kv("cc", "baseline_90d")
    cc_n    = _kv("cc", "n", 0) or 0
    cc_p    = _kv("cc", "prior")

    # Build four-lenses bullet list from nested findings
    fsig  = findings.get("failure_signatures", {})
    inter = findings.get("intersections", {})

    otd_sig = fsig.get("otd", {})
    top_reasons = otd_sig.get("top_reasons", [])
    top_reason_txt = top_reasons[0]["reason"] if top_reasons else "Unknown / not captured"

    ifr_sig  = fsig.get("ifr", {})
    alloc_pct = ifr_sig.get("allocation_gap_pct", 0) or 0
    hs_pct    = ifr_sig.get("hard_stockout_pct", 0) or 0
    cov_pct   = ifr_sig.get("covered_miss_pct", 0) or 0

    pm = inter.get("pfep_match", {})
    match_rate = pm.get("match_rate", 0) or 0

    triple_count = inter.get("triple", {}).get("triple_count", 0) or 0

    four_lenses = [
        f"Failure signature — top OTD reason: {top_reason_txt}",
        f"Allocation vs. stockout — alloc-gap {alloc_pct:.1f}%  |  hard-stockout {hs_pct:.1f}%",
        f"PFEP match rate on miss parts — {match_rate*100:.1f}%",
        f"Triple intersection (CC∩IFR∩OTD) — {triple_count} parts",
    ]

    # Realizations from findings
    raw_real = findings.get("realizations", []) or []
    realizations = []
    for r in raw_real[:4]:
        if isinstance(r, dict):
            realizations.append(r.get("realization", str(r))[:200])
        else:
            realizations.append(str(r)[:200])

    brain_insights = _generate_insights(findings)
    top_actions    = _top_actions(findings)

    now_str = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    date_range = f"{win_lo} → {win_hi}" if win_lo and win_hi else (anchor or "—")

    # ── Zone 0: Background ────────────────────────────────────────────────────
    _box(slide, 0, 0, SLIDE_W, SLIDE_H, fill="bg")

    # ── Zone 1: Header bar ────────────────────────────────────────────────────
    _box(slide, 0, 0, SLIDE_W, 0.72, fill="navy")
    # Title
    _text(slide, 0.25, 0.08, 8.0, 0.30,
          "Bi-Weekly Supply Chain Review",
          size=16, bold=True, color="bg")
    # Sub: site + date range
    _text(slide, 0.25, 0.40, 7.0, 0.25,
          f"Site: {site}   ·   Period: {date_range}",
          size=9, color="muted",
          align=PP_ALIGN.LEFT)
    # Timestamp right-aligned
    _text(slide, 9.50, 0.40, 3.60, 0.25,
          f"Generated: {now_str}",
          size=7, color="muted", align=PP_ALIGN.RIGHT)

    # ── Zone 2: KPI tiles (3 × equal width) ──────────────────────────────────
    TILE_T = 0.82
    TILE_H = 1.45
    TILE_W = (SLIDE_W - 0.50) / 3.0  # ~4.28"
    TILE_L = [0.25, 0.25 + TILE_W + 0.06, 0.25 + 2 * (TILE_W + 0.06)]

    _kpi_tile(slide, TILE_L[0], TILE_T, TILE_W, TILE_H,
              label="OTD 14D",
              value=_fmt(otd_14),
              delta=_delta_str(otd_14, otd_p),
              sub=f"goal 95.0%  |  90d {_fmt(otd_90)}  |  n={otd_n:,}",
              status=_kpi_status(otd_14))

    _kpi_tile(slide, TILE_L[1], TILE_T, TILE_W, TILE_H,
              label="IFR 14D",
              value=_fmt(ifr_14),
              delta=_delta_str(ifr_14, ifr_p),
              sub=f"goal 95.0%  |  90d {_fmt(ifr_90)}  |  n={ifr_n:,}",
              status=_kpi_status(ifr_14))

    _kpi_tile(slide, TILE_L[2], TILE_T, TILE_W, TILE_H,
              label="CYCLE-COUNT 14D",
              value=_fmt(cc_14),
              delta=_delta_str(cc_14, cc_p),
              sub=f"goal 95.0%  |  90d {_fmt(cc_90)}  |  n={cc_n:,}",
              status=_kpi_status(cc_14))

    # ── Zone 3: Four Lenses (left column) ────────────────────────────────────
    BODY_T = 2.40
    BODY_H = 3.90
    LEFT_W = 5.90
    RIGHT_L = 0.25 + LEFT_W + 0.15

    lens_texts = [str(l) for l in four_lenses[:6]] if four_lenses else [
        "No lens data available for this window."
    ]
    _bullet_block(slide, 0.25, BODY_T, LEFT_W, BODY_H / 2.0 - 0.05,
                  lens_texts, size=10, label="Four Lenses",
                  fill="insight", label_color="accent")

    # ── Zone 4: Realizations (left column, lower half) ───────────────────────
    real_texts = [str(r)[:160] for r in realizations[:4]] if realizations else [
        "No realizations for this window."
    ]
    _bullet_block(slide, 0.25, BODY_T + BODY_H / 2.0 + 0.05,
                  LEFT_W, BODY_H / 2.0 - 0.10,
                  real_texts, size=9, label="Realizations",
                  fill="tile_bg", label_color="accent")

    # ── Zone 5: Brain Insights (right column, upper) ─────────────────────────
    RIGHT_W = SLIDE_W - RIGHT_L - 0.25
    INSIGHT_H = BODY_H * 0.52
    if brain_insights:
        _bullet_block(slide, RIGHT_L, BODY_T, RIGHT_W, INSIGHT_H,
                      brain_insights, size=9.5,
                      label="Brain Insights",
                      fill="insight", label_color="accent",
                      bullet_color="ink")
    else:
        _bullet_block(slide, RIGHT_L, BODY_T, RIGHT_W, INSIGHT_H,
                      ["Run the analyzer to generate insights."], size=9.5,
                      label="Brain Insights", fill="insight", label_color="accent")

    # ── Zone 6: 30-Day Actions (right column, lower) ──────────────────────────
    ACTION_T = BODY_T + INSIGHT_H + 0.08
    ACTION_H = BODY_H - INSIGHT_H - 0.08
    action_bullets = [f"[T+30]  {a}" for a in top_actions]
    _bullet_block(slide, RIGHT_L, ACTION_T, RIGHT_W, ACTION_H,
                  action_bullets, size=9,
                  label="30-Day Actions",
                  fill="tile_bg", label_color="warn",
                  bullet_color="ink")

    # ── Zone 7: Footer ────────────────────────────────────────────────────────
    _box(slide, 0, 6.92, SLIDE_W, 0.08, fill="navy")
    seed = scope.get("seed", 9)
    policy = scope.get("anchor_policy", "AST-INV-PRO-0001")
    _text(slide, 0.25, 6.95, 12.80, 0.22,
          f"{policy}   ·   Site: {site}   ·   "
          f"Window: {date_range}   ·   Anchor: {anchor}   ·   seed {seed}   ·   {now_str}",
          size=7.5, color="muted", align=PP_ALIGN.LEFT)

    # ── Write ─────────────────────────────────────────────────────────────────
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out
