"""
One-Pager — single-slide 8.5×11 portrait Executive summary for a Mission.

Public API:
    render_one_pager(mission, result, viz, schema, out_path) -> Path

This is a *living* document: callers overwrite the same path on every
mission refresh so downstream consumers (agent_uplink, OneDrive viewers)
always see the latest version.
"""
from __future__ import annotations

import io
from datetime import datetime, timezone
from pathlib import Path

from . import builder as B


def _require():
    B._require_pptx()


def _fig_to_png_bytes(fig) -> bytes | None:
    """Render a Plotly figure to PNG. Returns None if kaleido is missing."""
    if fig is None:
        return None
    try:
        return fig.to_image(format="png", width=900, height=540, scale=2)
    except Exception:
        return None


def render_one_pager(mission, result, viz: dict, schema, out_path) -> Path:
    """Render the Executive 1-Pager (portrait, 8.5×11)."""
    _require()
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    # Initialize palette (also wires C_BG... globals used by builder helpers).
    p = B._palette()
    B.C_BG, B.C_INK, B.C_MUTED = p["BG"], p["INK"], p["MUTED"]
    B.C_ACCENT, B.C_GOOD, B.C_WARN, B.C_BAD = (
        p["ACCENT"], p["GOOD"], p["WARN"], p["BAD"]
    )

    prs = Presentation()
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11.0)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Quest banner
    quest_label = ""
    try:
        from ..brain import quests
        q = quests.get_quest(mission.quest_id)
        quest_label = q.label if q else mission.quest_id
    except Exception:
        quest_label = getattr(mission, "quest_id", "")

    B._text(slide, Inches(0.5), Inches(0.4), Inches(7.5), Inches(0.3),
            "Brain-Driven Quest · Executive 1-Pager",
            size=10, color=B.C_MUTED)
    B._text(slide, Inches(0.5), Inches(0.62), Inches(7.5), Inches(0.55),
            quest_label, size=22, bold=True, color=B.C_INK)
    B._hrule(slide, Inches(0.5), Inches(1.22), Inches(7.5))

    # Mission line
    site = getattr(mission, "site", "—") or "—"
    target_kind = getattr(mission, "target_entity_kind", "") or "—"
    target_key = getattr(mission, "target_entity_key", "") or "—"
    horizon = getattr(mission, "horizon_days", None) or 90
    summary = (
        f"Site: {site}   ·   Target: {target_kind}={target_key}   ·   "
        f"Horizon: {horizon}d   ·   Mission: {getattr(mission, 'id', '')}"
    )
    B._text(slide, Inches(0.5), Inches(1.35), Inches(7.5), Inches(0.3),
            summary, size=10, color=B.C_MUTED)

    # User query block
    user_q = getattr(mission, "user_query", "") or ""
    if user_q:
        B._text(slide, Inches(0.5), Inches(1.7), Inches(7.5), Inches(0.3),
                "USER QUERY", size=9, bold=True, color=B.C_MUTED)
        B._text(slide, Inches(0.5), Inches(1.95), Inches(7.5), Inches(0.7),
                user_q[:400], size=11, color=B.C_INK)

    # KPI tiles — top three from snapshot
    snap = getattr(result, "kpi_snapshot", {}) or {}
    kpis = sorted(((k, v) for k, v in snap.items()
                   if isinstance(v, (int, float))),
                  key=lambda x: -abs(float(x[1])))[:3]
    tile_top = Inches(2.85)
    tile_w = Inches(2.4)
    for i, (k, v) in enumerate(kpis):
        left = Inches(0.5 + i * 2.55)
        B._box(slide, left, tile_top, tile_w, Inches(1.1),
               fill=B.C_BG, line=B.C_MUTED)
        B._text(slide, left + Inches(0.15), tile_top + Inches(0.1),
                tile_w - Inches(0.3), Inches(0.3),
                k.upper().replace("_", " "),
                size=8, bold=True, color=B.C_MUTED)
        try:
            shown = f"{float(v):,.2f}"
        except Exception:
            shown = str(v)
        B._text(slide, left + Inches(0.15), tile_top + Inches(0.4),
                tile_w - Inches(0.3), Inches(0.6),
                shown, size=20, bold=True, color=B.C_INK)

    # Hero viz — pick the highest-priority figure that's available
    hero = (viz.get("kpi_trend") or viz.get("pareto") or
            viz.get("heatmap_matrix") or viz.get("network"))
    img_top = Inches(4.15)
    if hero is not None:
        png = _fig_to_png_bytes(hero)
        if png:
            slide.shapes.add_picture(io.BytesIO(png),
                                     Inches(0.5), img_top,
                                     width=Inches(7.5), height=Inches(3.6))
        else:
            # Fallback when kaleido missing — render caption + bullet summary.
            B._box(slide, Inches(0.5), img_top, Inches(7.5), Inches(3.6),
                   fill=B.C_BG, line=B.C_MUTED)
            B._text(slide, Inches(0.7), img_top + Inches(0.2),
                    Inches(7.1), Inches(0.4),
                    "Plotly figure (PNG export unavailable — install kaleido)",
                    size=11, bold=True, color=B.C_MUTED)
            from . import builder  # noqa: F401
            from .builder import _bullets
            try:
                from ..brain.viz_composer import caption_for
                cap = caption_for(hero)
            except Exception:
                cap = ""
            _bullets(slide, Inches(0.7), img_top + Inches(0.7),
                     Inches(7.1), Inches(2.7), [cap] if cap else [])

    # Top recommendations (from findings flagged with kind=recommendation)
    recs = []
    for f in (getattr(result, "findings", []) or []):
        if str(f.get("kind", "")).startswith("recommend"):
            recs.append(str(f.get("key", "")))
    if not recs:
        # Fall back to top-3 highest-score findings as recommendations.
        for f in sorted(getattr(result, "findings", []) or [],
                        key=lambda r: -float(r.get("score") or 0))[:3]:
            recs.append(str(f.get("key", "")))
    recs = recs[:3]

    rec_top = Inches(7.95)
    B._text(slide, Inches(0.5), rec_top, Inches(7.5), Inches(0.3),
            "TOP RECOMMENDATIONS", size=9, bold=True, color=B.C_MUTED)
    if recs:
        B._bullets(slide, Inches(0.5), rec_top + Inches(0.25),
                   Inches(7.5), Inches(1.4), recs, size=11)
    else:
        B._text(slide, Inches(0.5), rec_top + Inches(0.25), Inches(7.5),
                Inches(0.4), "No recommendations yet — refresh the mission.",
                size=11, color=B.C_MUTED)

    # Owner + progress bar + refresh footer
    foot_top = Inches(9.75)
    progress_raw = float(getattr(result, "progress_pct", 0.0) or 0.0)
    # Orchestrator emits 0..100; clamp + normalize for the bar.
    progress = max(0.0, min(1.0, progress_raw / 100.0))
    parsed = getattr(mission, "parsed_intent", {}) or {}
    owner = parsed.get("owner_role") or "Anyone"

    B._text(slide, Inches(0.5), foot_top, Inches(3.5), Inches(0.3),
            f"OWNER:  {owner}", size=10, bold=True, color=B.C_INK)
    refreshed = (getattr(mission, "last_refreshed_at", None)
                 or datetime.now(timezone.utc).isoformat())
    B._text(slide, Inches(4.0), foot_top, Inches(4.0), Inches(0.3),
            f"Refreshed: {refreshed[:19].replace('T', ' ')} UTC",
            size=9, color=B.C_MUTED, align=2)  # PP_ALIGN.RIGHT == 2

    # Progress bar
    bar_top = foot_top + Inches(0.4)
    B._box(slide, Inches(0.5), bar_top, Inches(7.5), Inches(0.22),
           fill=B.C_BG, line=B.C_MUTED)
    if progress > 0:
        B._box(slide, Inches(0.5), bar_top, Inches(7.5 * progress),
               Inches(0.22), fill=B.C_ACCENT, line=B.C_ACCENT)
    B._text(slide, Inches(0.5), bar_top + Inches(0.3), Inches(7.5),
            Inches(0.3), f"Progress: {progress*100:.0f}%",
            size=9, color=B.C_MUTED)
    # Save (overwrite-in-place for living doc)
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


__all__ = ["render_one_pager"]
