"""
Implementation Plan PPTX — multi-slide living document for a Mission.

Slide order (fixed):
    1. Cover
    2. Mission & Quest context
    3. Entity schema (Mermaid as PNG when CLI present, else table fallback)
    4. Current state — KPI snapshot
    5. Root-cause findings (top by score)
    6. Recommended interventions (systemic vs operational)
    7. Phased rollout
    8. Dependencies & risks
    9. Appendix — analyzer dispatch + raw scope tags
"""
from __future__ import annotations

import io
import json
import shutil
import subprocess
import tempfile
from datetime import datetime, timezone
from pathlib import Path

from . import builder as B


def _require():
    B._require_pptx()


def _palette_init():
    p = B._palette()
    B.C_BG, B.C_INK, B.C_MUTED = p["BG"], p["INK"], p["MUTED"]
    B.C_ACCENT, B.C_GOOD, B.C_WARN, B.C_BAD = (
        p["ACCENT"], p["GOOD"], p["WARN"], p["BAD"]
    )


def _fig_to_png(fig) -> bytes | None:
    if fig is None:
        return None
    try:
        return fig.to_image(format="png", width=1100, height=560, scale=2)
    except Exception:
        return None


def _mermaid_to_png(mermaid: str) -> bytes | None:
    """Render Mermaid via mmdc CLI if available, else return None."""
    if not mermaid or shutil.which("mmdc") is None:
        return None
    try:
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "er.mmd"
            dst = Path(td) / "er.png"
            src.write_text(mermaid, encoding="utf-8")
            subprocess.run(
                ["mmdc", "-i", str(src), "-o", str(dst), "-b", "white"],
                capture_output=True, check=True, timeout=30,
            )
            if dst.exists():
                return dst.read_bytes()
    except Exception:
        return None
    return None


def _slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def _s_cover(prs, mission, quest_label):
    from pptx.util import Inches
    s = _slide_blank(prs)
    B._box(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
           fill=B.C_BG, line=B.C_BG)
    B._text(s, Inches(0.7), Inches(2.4), Inches(11.5), Inches(0.4),
            "BRAIN-DRIVEN QUEST", size=12, bold=True, color=B.C_ACCENT)
    B._text(s, Inches(0.7), Inches(2.85), Inches(11.5), Inches(1.0),
            "Implementation Plan", size=42, bold=True, color=B.C_INK)
    B._text(s, Inches(0.7), Inches(3.95), Inches(11.5), Inches(0.5),
            quest_label, size=20, color=B.C_INK)
    B._text(s, Inches(0.7), Inches(4.55), Inches(11.5), Inches(0.4),
            f"Site: {mission.site}   ·   Mission: {mission.id}",
            size=14, color=B.C_MUTED)
    B._text(s, Inches(0.7), Inches(6.7), Inches(11.5), Inches(0.4),
            f"Generated {datetime.now(timezone.utc).isoformat()[:19]} UTC",
            size=10, color=B.C_MUTED)


def _s_context(prs, mission, result):
    from pptx.util import Inches
    s = _slide_blank(prs)
    B._header(s, "Mission & Quest Context", eyebrow="Implementation Plan")
    user_q = (mission.user_query or "").strip()
    parsed = mission.parsed_intent or {}
    lines = [
        f"User query: {user_q[:300]}",
        f"Quest:      {mission.quest_id}",
        f"Site:       {mission.site}",
        f"Target:     {mission.target_entity_kind} = {mission.target_entity_key}",
        f"Horizon:    {mission.horizon_days} days",
        f"Owner:      {parsed.get('owner_role', 'Anyone')}",
        f"Scope tags: {', '.join(mission.scope_tags) or '—'}",
        f"Parsed by:  {parsed.get('parser_source', 'unknown')}",
    ]
    B._bullets(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
               lines, size=14)


def _s_schema(prs, schema):
    from pptx.util import Inches
    s = _slide_blank(prs)
    B._header(s, "Entity Schema", eyebrow="Implementation Plan")

    png = _mermaid_to_png(getattr(schema, "mermaid", "") or "")
    if png:
        s.shapes.add_picture(io.BytesIO(png), Inches(0.5), Inches(1.4),
                             width=Inches(12.3), height=Inches(5.6))
        return

    # Table fallback per plan §"Further Considerations" #2
    headers = ["Logical table", "Qualified", "Cols mapped", "Sample"]
    rows = []
    for t in (schema.tables or [])[:14]:
        mapped = sum(1 for c in t.columns if c.logical)
        sample = ", ".join(c.logical or c.name for c in t.columns
                           if c.logical)[:60]
        rows.append([t.logical_name, t.qualified, f"{mapped}/{len(t.columns)}",
                     sample or "—"])
    if rows:
        B._table(s, Inches(0.5), Inches(1.5), Inches(12.3),
                 Inches(min(5.5, 0.4 + 0.35 * len(rows))), headers, rows)
    else:
        B._text(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.5),
                "No tables resolved for this entity kind.",
                size=14, color=B.C_MUTED)

    # Relationship summary at bottom
    rels = getattr(schema, "relationships", []) or []
    if rels:
        rel_lines = [f"{r.from_table} ⇄ {r.to_table}  on {r.on}"
                     for r in rels[:8]]
        B._text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.3),
                "RELATIONSHIPS", size=9, bold=True, color=B.C_MUTED)
        B._bullets(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(1.1),
                   rel_lines, size=10)


def _s_current_state(prs, viz):
    from pptx.util import Inches
    s = _slide_blank(prs)
    B._header(s, "Current State — KPI Snapshot", eyebrow="Implementation Plan")
    fig = viz.get("kpi_trend") or viz.get("pareto")
    png = _fig_to_png(fig)
    if png:
        s.shapes.add_picture(io.BytesIO(png), Inches(0.5), Inches(1.4),
                             width=Inches(12.3), height=Inches(5.6))
    else:
        B._text(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.5),
                "No KPIs available — analyzers returned empty.",
                size=14, color=B.C_MUTED)


def _s_findings(prs, result):
    from pptx.util import Inches
    s = _slide_blank(prs)
    B._header(s, "Root-Cause Findings", eyebrow="Implementation Plan")
    findings = sorted(getattr(result, "findings", []) or [],
                      key=lambda r: -float(r.get("score") or 0))[:12]
    if not findings:
        B._text(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.5),
                "No findings produced this round.",
                size=14, color=B.C_MUTED)
        return
    headers = ["Page/Module", "Kind", "Key", "Score"]
    rows = [[str(f.get("page", ""))[:24],
             str(f.get("kind", ""))[:24],
             str(f.get("key", ""))[:48],
             f"{float(f.get('score') or 0):.2f}"] for f in findings]
    B._table(s, Inches(0.5), Inches(1.5), Inches(12.3),
             Inches(min(5.7, 0.4 + 0.35 * len(rows))), headers, rows)


def _split_recommendations(result):
    """Return (systemic, operational) lists per CrossDataset spec §6.

    Heuristic: kind starts with 'recommend_systemic' / 'recommend_operational',
    else partition by score (>=0.7 systemic, else operational).
    """
    sys_, ops = [], []
    for f in (getattr(result, "findings", []) or []):
        kind = str(f.get("kind", ""))
        key = str(f.get("key", ""))
        if "systemic" in kind:
            sys_.append(key)
        elif "operational" in kind or "operation" in kind:
            ops.append(key)
        elif kind.startswith("recommend"):
            (sys_ if float(f.get("score") or 0) >= 0.7 else ops).append(key)
    if not (sys_ or ops):
        for f in sorted(getattr(result, "findings", []) or [],
                        key=lambda r: -float(r.get("score") or 0))[:6]:
            (sys_ if float(f.get("score") or 0) >= 0.7
             else ops).append(str(f.get("key", "")))
    return sys_[:6], ops[:6]


def _s_recommendations(prs, result):
    from pptx.util import Inches
    s = _slide_blank(prs)
    B._header(s, "Recommended Interventions", eyebrow="Implementation Plan")
    sys_, ops = _split_recommendations(result)
    B._text(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.4),
            "SYSTEMIC", size=12, bold=True, color=B.C_ACCENT)
    B._bullets(s, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0),
               sys_ or ["—"], size=12)
    B._text(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.4),
            "OPERATIONAL", size=12, bold=True, color=B.C_GOOD)
    B._bullets(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0),
               ops or ["—"], size=12)


def _s_rollout(prs, result):
    from pptx.util import Inches
    s = _slide_blank(prs)
    B._header(s, "Phased Rollout", eyebrow="Implementation Plan")
    sys_, ops = _split_recommendations(result)
    phases = [
        ("Phase 1 — Stabilize (Week 1-2)",
         (ops[:2] or ["Establish baseline KPIs", "Confirm data quality"])),
        ("Phase 2 — Implement (Week 3-6)",
         (sys_[:3] or ["Deploy systemic recommendations"])),
        ("Phase 3 — Sustain (Week 7+)",
         ["Mission auto-refresh continues",
          "Living artifacts overwrite in place",
          "Brain surfaces drift via Body directives"]),
    ]
    y = 1.4
    for header, items in phases:
        B._text(s, Inches(0.5), Inches(y), Inches(12.3), Inches(0.3),
                header, size=12, bold=True, color=B.C_INK)
        B._bullets(s, Inches(0.7), Inches(y + 0.35), Inches(12.0),
                   Inches(1.4), items, size=11)
        y += 1.85


def _s_risks(prs, result):
    from pptx.util import Inches
    s = _slide_blank(prs)
    B._header(s, "Dependencies & Risks", eyebrow="Implementation Plan")
    deps = ["edap_dw_replica connectivity",
            "Open-weight LLM ensemble availability",
            "python-pptx, plotly, kaleido (PNG export)",
            "mermaid-cli (mmdc) for ER diagram render — falls back to table"]
    risks = ["Analyzer outage degrades scope coverage (handled — _safe_run)",
             "LLM intent mis-parse (handled — keyword fallback)",
             "Stale baseline if first refresh fails — re-run mission",
             "Schema drift in DW — refresh discovered_schema.yaml"]
    B._text(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.3),
            "DEPENDENCIES", size=12, bold=True, color=B.C_ACCENT)
    B._bullets(s, Inches(0.5), Inches(1.75), Inches(6.0), Inches(5.0),
               deps, size=11)
    B._text(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.3),
            "RISKS", size=12, bold=True, color=B.C_BAD)
    B._bullets(s, Inches(6.8), Inches(1.75), Inches(6.0), Inches(5.0),
               risks, size=11)


def _s_appendix(prs, mission, result, viz):
    from pptx.util import Inches
    s = _slide_blank(prs)
    B._header(s, "Appendix — Dispatch & Diagnostics",
              eyebrow="Implementation Plan")
    fig = viz.get("network") or viz.get("sankey_flow")
    png = _fig_to_png(fig)
    if png:
        s.shapes.add_picture(io.BytesIO(png), Inches(0.5), Inches(1.4),
                             width=Inches(8.5), height=Inches(5.6))
    # Right-side text summary
    outcomes = getattr(result, "outcomes", []) or []
    lines = [f"{o.scope_tag} → {o.analyzer}: "
             f"{'OK' if o.ok else 'FAIL'} ({o.n_findings})"
             for o in outcomes[:14]]
    B._text(s, Inches(9.2), Inches(1.4), Inches(3.6), Inches(0.3),
            "ANALYZERS", size=10, bold=True, color=B.C_MUTED)
    B._bullets(s, Inches(9.2), Inches(1.7), Inches(3.6), Inches(5.0),
               lines or ["(no analyzer outcomes)"], size=9)
    B._text(s, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
            f"Mission {mission.id}  ·  refreshed "
            f"{datetime.now(timezone.utc).isoformat()[:19]} UTC  ·  "
            f"progress {float(getattr(result, 'progress_pct', 0.0) or 0):.0f}%",
            size=9, color=B.C_MUTED)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------
def render_implementation_plan(mission, result, viz: dict, schema, out_path) -> Path:
    """Render the multi-slide Implementation Plan PPTX (overwrite in place)."""
    _require()
    _palette_init()
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    quest_label = mission.quest_id
    try:
        from ..brain import quests
        q = quests.get_quest(mission.quest_id)
        if q:
            quest_label = q.label
    except Exception:
        pass

    _s_cover(prs, mission, quest_label)
    _s_context(prs, mission, result)
    _s_schema(prs, schema)
    _s_current_state(prs, viz)
    _s_findings(prs, result)
    _s_recommendations(prs, result)
    _s_rollout(prs, result)
    _s_risks(prs, result)
    _s_appendix(prs, mission, result, viz)

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


__all__ = ["render_implementation_plan"]
